from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation object
prs = Presentation()

# Slide titles and content based on the user's request
slides_content = [
    ("Introduction to Hadoop", "Understanding Big Data and Hadoop's Role\nPresenter: Richard Kirschner\nOrganization: Simply Learn\nWelcome message and brief introduction to the topic."),
    ("Welcome to Hadoop", "Overview of Hadoop and its relevance in managing big data.\nMention Simply Learn and their focus on providing certification.\nHighlight the importance of Hadoop in the modern data landscape.\nBriefly outline the structure of the presentation."),
    ("Understanding Hadoop through a Farm Analogy", "A visual analogy using a farm to explain distributed data processing.\nIntroduce the farm story as a way to simplify understanding of complex Hadoop concepts."),
    ("Jack’s Fruit Farm", "Jack harvests different fruits and stores them.\nJack harvests grapes, apples, and oranges.\nInitial storage in a single shed.\nThe challenge of handling increased variety and volume of fruits."),
    ("Hiring Help and Distributed Storage", "Jack hires more workers and distributes storage.\nHiring two additional workers to help with harvesting.\nDistributing the storage to multiple sheds to handle more fruits.\nExplanation of parallel processing in the context of the farm."),
    ("From Farm to Big Data", "How the farm analogy relates to big data processing.\nThe farm workers as processors.\nThe fruits as data.\nThe storage sheds as distributed storage units."),
    ("Evolution from Single Processor to Distributed Systems", "Historical context of data processing.\nEarly days of single processors and mainframes.\nTransition to modern distributed systems.\nImportance of handling increasing data volumes and complexity."),
    ("Structured, Semi-structured, and Unstructured Data", "Definitions and examples.\nStructured data: Databases, SQL queries.\nSemi-structured data: Emails, XML, HTML.\nUnstructured data: Photos, videos, social media posts."),
    ("Challenges with Big Data", "Volume, Velocity, Variety, Value, and Veracity.\nVolume: The sheer amount of data generated.\nVelocity: The speed at which data is generated.\nVariety: Different types and formats of data.\nValue: Extracting meaningful insights from data.\nVeracity: Ensuring data quality and accuracy."),
    ("Introducing Hadoop", "Hadoop as a solution to big data challenges.\nHadoop's ability to handle large volumes of data.\nDistributed storage and processing capabilities.\nFlexibility in handling various data types."),
    ("Hadoop Ecosystem", "Key components - HDFS, MapReduce, and YARN.\nHDFS: Hadoop Distributed File System.\nMapReduce: Parallel processing model.\nYARN: Yet Another Resource Negotiator, manages resources."),
    ("HDFS Overview", "Purpose and function of HDFS.\nDesigned for storing large data sets on commodity hardware.\nManages storage across multiple nodes."),
    ("HDFS Architecture", "NameNode and DataNodes.\nNameNode: Manages metadata and controls DataNodes.\nDataNodes: Store actual data, perform read/write operations."),
    ("Importance of Commodity Hardware", "Cost benefits of using commodity hardware.\nComparison with enterprise hardware.\nCost-effectiveness of using commodity servers for DataNodes.\nEnterprise hardware often used for NameNode."),
    ("Data Replication in HDFS", "How HDFS ensures data reliability.\nReplication of data blocks across multiple DataNodes.\nDefault replication factor of three.\nEnsures fault tolerance and data availability."),
    ("Introduction to MapReduce", "MapReduce programming model.\nParallel processing framework.\nDivides tasks into Map and Reduce functions."),
    ("MapReduce Workflow", "Steps in MapReduce - Map, Shuffle, Reduce.\nMap: Processes input data and generates key-value pairs.\nShuffle: Organizes data by keys.\nReduce: Aggregates and summarizes the data."),
    ("MapReduce Example", "Word count example.\nInput: Text data.\nMap: Count occurrences of words.\nShuffle and Sort: Organize word counts.\nReduce: Sum the counts for each word."),
    ("What is YARN?", "YARN as the resource management layer.\nManages resources and job scheduling.\nCoordinates tasks across the Hadoop cluster."),
    ("YARN Components", "ResourceManager, NodeManager, and ApplicationMaster.\nResourceManager: Oversees resource allocation.\nNodeManager: Manages individual nodes.\nApplicationMaster: Manages application execution."),
    ("Hadoop Use Case", "Example of Hadoop in combating fraud.\nOverview of Zions Bank use case.\nApplication of Hadoop in real-world scenarios."),
    ("Zions Bank Challenges", "Issues with storing and analyzing large data.\nLimitations of traditional RDBMS.\nNeed for handling unstructured and semi-structured data."),
    ("Hadoop Implementation", "Steps taken to implement Hadoop.\nAdoption of Hadoop for data storage and processing.\nUse of HDFS and MapReduce."),
    ("Benefits of Hadoop", "Improved fraud detection and data analysis.\nEnhanced ability to analyze large datasets.\nIncreased efficiency and accuracy in detecting fraud."),
    ("Summary of Key Points", "Review of big data challenges and Hadoop solutions.\nImportance of distributed storage and processing.\nHadoop's role in modern data management."),
    ("Certification Information", "Benefits of Hadoop certification.\nHow to get certified with Simply Learn.\nAdvantages of certification for career growth."),
    ("Conclusion", "Final thoughts and closing remarks.\nEncouragement to explore Hadoop further.\nSummary of the importance of Hadoop in big data."),
    ("Learning Resources", "Links to further reading and tutorials.\nSimply Learn resources.\nOther recommended materials for in-depth learning."),
    ("Contact Us", "Simply Learn contact details.\nHow to reach out for more information.\nEmail and website details."),
    ("Questions & Answers", "Open floor for questions.\nEncourage audience participation.\nAddress common queries about Hadoop."),
    ("HDFS Architecture Diagram", "Visual representation of NameNode and DataNodes.\nDiagram showing the interaction between NameNode and DataNodes.\nLabels and annotations explaining the components."),
    ("MapReduce Workflow Diagram", "Diagram of Map, Shuffle, and Reduce phases.\nDetailed process flow.\nVisual representation of data flow through MapReduce phases."),
    ("YARN Architecture Diagram", "Visual representation of ResourceManager and NodeManagers.\nLabels and annotations explaining the components.\nVisual explanation of resource management and job scheduling."),
    ("Fraud Detection Process", "Diagram showing data flow in fraud detection.\nSteps and tools used in the process.\nVisual representation of how Hadoop aids in fraud detection."),
    ("How to Get Certified", "Steps to achieve Hadoop certification.\nRequirements and preparation tips.\nBenefits of certification for professionals."),
    ("Big Data Challenges Recap", "Summary of Volume, Velocity, Variety, Value, and Veracity.\nBrief explanations of each challenge.\nImportance of addressing these challenges."),
    ("Hadoop Components Recap", "Summary of HDFS, MapReduce, and YARN.\nKey functions of each component.\nImportance in the Hadoop ecosystem."),
    ("Use Case Recap", "Summary of Zions Bank case study.\nKey steps in implementing Hadoop.\nBenefits realized by Zions Bank."),
    ("Hadoop Certification Recap", "Summary of certification process.\nSteps to achieve certification.\nBenefits for career advancement."),
    ("Final Thoughts", "Encouragement and motivation to learn Hadoop.\nEmphasize the importance of Hadoop in the current data landscape.\nInspire the audience to pursue further learning and certification."),
    ("What is Big Data?", "Defining Big Data and its significance.\nExplanation of Big Data characteristics.\nImportance in modern data analysis and decision-making."),
    ("Key Characteristics of Big Data", "In-depth look at Volume, Velocity, Variety, Value, and Veracity.\nDetailed examples and implications of each characteristic."),
    ("Examples of Big Data", "Real-world examples of Big Data applications.\nSocial media analytics, healthcare data, financial transactions."),
    ("Limitations of Traditional Data Processing", "Challenges faced with traditional systems.\nInability to handle large volumes and diverse data types.\nPerformance bottlenecks and scalability issues."),
    ("Detailed Components of Hadoop Ecosystem", "Overview of additional Hadoop ecosystem tools.\nApache Pig, Apache Hive, Apache HBase, and Apache Spark."),
    ("Introduction to Apache Pig", "High-level platform for creating MapReduce programs.\nBenefits and use cases of Apache Pig.\nExample of a Pig Latin script."),
    ("Introduction to Apache Hive", "Data warehouse infrastructure on top of Hadoop.\nBenefits and use cases of Apache Hive.\nExample of a Hive query."),
    ("Introduction to Apache HBase", "NoSQL database for Hadoop.\nBenefits and use cases of Apache HBase.\nExample of HBase usage."),
    ("Introduction to Apache Spark", "Fast and general-purpose cluster computing system.\nBenefits and use cases of Apache Spark.\nExample of a Spark application."),
    ("How HDFS Stores Data", "Detailed process of data storage in HDFS.\nData blocks, replication, and fault tolerance.\nExample of data distribution across nodes."),
    ("HDFS Read and Write Operations", "How data is read from and written to HDFS.\nSteps involved in read and write operations.\nEnsuring data consistency and reliability."),
    ("MapReduce Detailed Workflow", "In-depth explanation of MapReduce phases.\nMap phase, Shuffle phase, and Reduce phase.\nExamples and visual aids."),
    ("Optimizing MapReduce Jobs", "Techniques to improve MapReduce performance.\nCombiner functions, data locality, and partitioning."),
    ("YARN Resource Management", "How YARN manages cluster resources.\nAllocation of CPU, memory, and storage.\nScheduling and monitoring tasks."),
    ("Managing Applications with YARN", "How YARN handles application requests.\nResource allocation process.\nApplication lifecycle management."),
    ("Hadoop Security", "Ensuring data security in Hadoop.\nAuthentication, authorization, and encryption.\nSecurity best practices."),
    ("Hadoop Performance Tuning", "Techniques for optimizing Hadoop performance.\nHardware configurations, cluster management, and job optimization."),
    ("Hadoop Use Cases", "Examples of industries using Hadoop.\nHealthcare, finance, retail, and social media."),
    ("Hadoop Community and Support", "Resources for learning and troubleshooting Hadoop.\nOnline forums, documentation, and support channels."),
    ("Future of Hadoop", "The Future of Hadoop\nEmerging trends and advancements in Hadoop.\nIntegration with cloud computing, AI, and machine learning.")
]

# Add slides to the presentation
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Use the 'Title and Content' layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

# Save the presentation
pptx_path = "Hadoop_Presentation.pptx"
prs.save(pptx_path)